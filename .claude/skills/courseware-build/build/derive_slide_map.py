#!/usr/bin/env python3
"""Derive slide anchors from the BUILT deck and write slide_map.py.

The Lesson Plan cites slide numbers. Hand-maintained numbers drift the moment a
slide is added or removed, which sends a trainer to the wrong slide at every
transition. This script reads the actual .pptx and emits the real numbers, so
build_lesson_plan.py can never be wrong. It runs as part of the orchestrator,
after the deck is built and before the LP.
"""
import os, re, sys, glob
from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__)); sys.path.insert(0, HERE)
import course_data as C


def _find_repo(start):
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env): return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "activities")):
            return d
    return os.path.dirname(os.path.dirname(HERE))


REPO = _find_repo(HERE)
decks = sorted(glob.glob(os.path.join(REPO, "courseware", "*.pptx")), key=os.path.getmtime)
if not decks:
    raise SystemExit("no deck found — build the slides first")
prs = Presentation(decks[-1])


def text(s):
    return " ".join(" ".join(sh.text_frame.text.split()) for sh in s.shapes
                    if sh.has_text_frame and sh.text_frame.text.strip())


T = [text(s) for s in prs.slides]          # T[i] is slide i+1
N = len(T)


def first(pred, lo=1, hi=None):
    hi = hi or N
    for i in range(lo, hi + 1):
        if pred(T[i - 1]):
            return i
    return None


def all_hits(pred):
    return [i for i in range(1, N + 1) if pred(T[i - 1])]


# ---- activity blocks: the ACTIVITY n briefing slide through its debrief slide ----
acts = {}
brief = {}
for i, t in enumerate(T, 1):
    m = re.search(r"\bACTIVITY (\d+)\b", t)
    if m and "HANDS-ON" in t:
        brief.setdefault(int(m.group(1)), i)
for n, s in brief.items():
    # the block ends at THIS activity's debrief slide (kicker "ACTIVITY n · WHAT WE LEARNED")
    tag = "ACTIVITY %d · WHAT WE LEARNED" % n
    e = first(lambda x, tag=tag: tag in x.replace("  ", " "), lo=s)
    if e is None:
        # fall back: end just before the next activity's briefing
        nxt = min((v for k, v in brief.items() if v > s), default=N)
        e = nxt - 1
    acts[n] = (s, e)

# ---- section dividers: a divider carries the big topic number and no kicker body ----
def divider(num, title_frag):
    return first(lambda t: t.startswith("TOPIC %d" % num) or
                 ("TOPIC %d" % num in t[:40] and title_frag.lower() in t.lower()))


topics = {}
for tp in C.TOPICS:
    n = tp["num"]
    frag = tp["title"].split()[0]
    # the divider slide shows "TOPIC n" as its kicker AND the topic title AND the big numeral
    hit = first(lambda t, n=n, tp=tp: ("TOPIC %d" % n) in t and tp["title"] in t
                and tp["subtitle"][:24] in t)
    topics[n] = hit or divider(n, frag)

ANCHORS = dict(
    total=N,
    activities=acts,
    topics=topics,
    admin_start=1,
    ice_breaker=first(lambda t: "Let's know each other" in t),
    trainer_template=first(lambda t: "GENERAL TRAINER TEMPLATE" in t),
    trainer_named=first(lambda t: C.TRAINER in t and "YOUR TRAINER" in t),
    ground_rules=first(lambda t: "Ground Rules" in t),
    download=first(lambda t: "Download Your Course Material" in t),
    skills_framework=first(lambda t: "Skills Framework Alignment" in t),
    outcomes=first(lambda t: "Learning Outcomes" in t and "BY THE END" in t),
    lesson_plan=all_hits(lambda t: "Lesson Plan — Day" in t),
    briefing=all_hits(lambda t: "Briefing for Assessment" in t),
    assessment=all_hits(lambda t: "Final Assessment" in t and "Written Assessment (WA)" in t),
    flow=all_hits(lambda t: "Assessment Flow" in t),
    attendance=all_hits(lambda t: "Digital Attendance" in t and "TRAQOM" in t),
    # NB: the Day 1/Day 2 lesson-plan TABLE slides quote "Lunch Break" and
    # "End of Day 1" as row text, so a naive substring match lands on the table
    # rather than the real divider. Divider slides are short and start with the
    # phrase, so anchor on that.
    lunch=all_hits(lambda t: t.strip().startswith("Lunch Break")),
    end_day1=first(lambda t: t.strip().startswith("End of Day 1")),
    recap_day1=first(lambda t: "Day 1 in one line" in t),
    toolkit=first(lambda t: "Four tools. Four different jobs." in t
                  or "Choosing Your Root Cause Tool" in t),
    convergence=first(lambda t: "Divergence Then Convergence" in t),
    implementation=first(lambda t: "From Decision to Proven Result" in t),
    summary=first(lambda t: "The Complete Method" in t),
    traqom_cert=first(lambda t: "Certificate & TRAQOM" in t),
    thanks=first(lambda t: "Thank You" in t),
)

OUT = os.path.join(HERE, "slide_map.py")
with open(OUT, "w") as f:
    f.write('"""AUTO-GENERATED from the built deck by derive_slide_map.py.\n'
            'DO NOT HAND-EDIT — regenerated on every build so the Lesson Plan\n'
            'slide references can never drift from the deck."""\n\n')
    f.write("ANCHORS = %r\n" % (ANCHORS,))

print("slide_map.py written from:", os.path.basename(decks[-1]))
print("  slides:", N)
print("  activities:", {k: acts[k] for k in sorted(acts)})
print("  topics:", topics)
print("  briefing:", ANCHORS["briefing"], " assessment:", ANCHORS["assessment"],
      " flow:", ANCHORS["flow"], " attendance:", ANCHORS["attendance"])
missing = [k for k, v in ANCHORS.items() if v is None]
if missing:
    raise SystemExit("ANCHOR NOT FOUND: %s — the deck changed; fix the matcher." % missing)

# ---- sanity assertions -------------------------------------------------------
# A missing anchor is caught above. These catch the nastier case: an anchor that
# is PRESENT but points at the wrong slide (e.g. matching the lesson-plan table's
# row text instead of the real divider), which silently produces wrong LP refs.
t1, t2, t3 = topics[1], topics[2], topics[3]
problems = []
if not (t1 < t2 < t3):
    problems.append("topic dividers out of order: %s" % topics)
if ANCHORS["end_day1"] <= t1:
    problems.append("end_day1 (%d) must come after the Topic 1 divider (%d)"
                    % (ANCHORS["end_day1"], t1))
if not (ANCHORS["end_day1"] < t2):
    problems.append("end_day1 (%d) must precede the Topic 2 divider (%d)"
                    % (ANCHORS["end_day1"], t2))
if ANCHORS["recap_day1"] is not None and not (ANCHORS["end_day1"] < ANCHORS["recap_day1"] < t2):
    problems.append("Day 1 recap (%s) must sit between End of Day 1 (%d) and Topic 2 (%d)"
                    % (ANCHORS["recap_day1"], ANCHORS["end_day1"], t2))
early_lunch = [s for s in ANCHORS["lunch"] if s < t1]
if early_lunch:
    problems.append("lunch anchors before the Topic 1 divider: %s" % early_lunch)
for n, (lo, hi) in acts.items():
    if lo > hi:
        problems.append("activity %d range inverted: %d-%d" % (n, lo, hi))
    if n > 1 and lo <= acts[n - 1][1]:
        problems.append("activity %d (%d) overlaps activity %d (ends %d)"
                        % (n, lo, n - 1, acts[n - 1][1]))
if sorted(acts) != list(range(1, len(acts) + 1)):
    problems.append("activity numbering not contiguous: %s" % sorted(acts))
if ANCHORS["briefing"][-1] >= ANCHORS["assessment"][-1]:
    problems.append("closing Briefing (%d) must precede Assessment (%d)"
                    % (ANCHORS["briefing"][-1], ANCHORS["assessment"][-1]))
if problems:
    raise SystemExit("ANCHOR SANITY FAILED:\n  - " + "\n  - ".join(problems))
print("  anchor sanity checks: OK")
